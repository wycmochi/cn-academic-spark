from __future__ import annotations

import io

import matplotlib
import matplotlib.pyplot as plt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

from CN_Spark_paper2ppt.scripts.layout_library import (
    FONTS,
    LAYOUT,
    THEME,
    add_bottom_banner,
    add_citation_footer,
    add_header,
    rgb,
)

matplotlib.rcParams["font.family"] = ["SimHei", "DejaVu Sans"]
matplotlib.rcParams["axes.unicode_minus"] = False


def formula_to_png(latex_str: str, fontsize=14, width=5.2, height=0.95) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(width, height))
    ax.text(0.5, 0.5, f"${latex_str}$", fontsize=fontsize, ha="center", va="center", transform=ax.transAxes)
    ax.axis("off")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight", transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf


def insert_policy_image(slide, image_path_or_bytes, x_cm, y_cm, w_cm, h_cm, caption_text: str, theme, fonts):
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(x_cm - 0.05),
        Cm(y_cm - 0.05),
        Cm(w_cm + 0.1),
        Cm(h_cm + 0.35),
    )
    border.fill.solid()
    border.fill.fore_color.rgb = rgb("F0F4FA")
    border.line.color.rgb = rgb(theme["primary"])
    border.line.width = Pt(1.2)

    slide.shapes.add_picture(image_path_or_bytes, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    cap_box = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm + h_cm + 0.05), Cm(w_cm), Cm(0.3))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption_text
    run.font.name = "微软雅黑"
    run.font.size = Pt(9)
    run.font.color.rgb = rgb(fonts["caption"][3])


def insert_policy_stat_card(slide, stat_number: str, stat_label: str, source_text: str, x_cm, y_cm, w_cm, h_cm, theme):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(theme["primary"])
    card.line.fill.background()

    num_box = slide.shapes.add_textbox(Cm(x_cm + 0.2), Cm(y_cm + 0.22), Cm(w_cm - 0.4), Cm(h_cm * 0.42))
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.alignment = PP_ALIGN.CENTER
    run = num_p.add_run()
    run.text = stat_number
    run.font.name = "Times New Roman"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = rgb("FFC000")

    label_box = slide.shapes.add_textbox(Cm(x_cm + 0.2), Cm(y_cm + h_cm * 0.47), Cm(w_cm - 0.4), Cm(h_cm * 0.28))
    label_tf = label_box.text_frame
    label_p = label_tf.paragraphs[0]
    label_p.alignment = PP_ALIGN.CENTER
    run = label_p.add_run()
    run.text = stat_label
    run.font.name = "微软雅黑"
    run.font.size = Pt(13)
    run.font.color.rgb = rgb("FFFFFF")

    src_box = slide.shapes.add_textbox(Cm(x_cm + 0.2), Cm(y_cm + h_cm - 0.35), Cm(w_cm - 0.4), Cm(0.25))
    src_tf = src_box.text_frame
    src_p = src_tf.paragraphs[0]
    src_p.alignment = PP_ALIGN.CENTER
    run = src_p.add_run()
    run.text = f"来源：{source_text}"
    run.font.name = "微软雅黑"
    run.font.size = Pt(8)
    run.font.color.rgb = rgb("AAAAAA")


def generate_gantt_chart(gantt_data: dict, theme: dict) -> io.BytesIO:
    tasks = gantt_data["tasks"]
    total_months = gantt_data["total_months"]
    n_tasks = len(tasks)

    fig, ax = plt.subplots(figsize=(14, max(4, n_tasks * 0.55 + 1.5)))
    fig.patch.set_facecolor("#F8F9FA")
    ax.set_facecolor("#FFFFFF")

    task_names = [t["name"] for t in tasks]
    y_positions = list(range(n_tasks, 0, -1))

    for task, y in zip(tasks, y_positions):
        start = task["start"] - 1
        duration = task["duration"]
        color = f"#{task['color']}"
        if task.get("milestone"):
            ax.plot(start + duration, y, "D", markersize=12, color=color, zorder=5)
        else:
            ax.barh(y, duration, left=start, height=0.5, color=color, alpha=0.88, edgecolor="white", linewidth=0.8)
            ax.text(start + duration / 2, y, task["name"], ha="center", va="center", fontsize=8.5, color="white", fontweight="bold", zorder=6)

    ax.set_xlim(0, total_months)
    ax.set_xticks(range(total_months + 1))
    start_label = gantt_data.get("start_label", "2025年09月")
    start_year = int(start_label[:4])
    start_month = int(start_label[5:7]) if "月" in start_label else 9
    month_labels = []
    for m in range(total_months + 1):
        month = (start_month - 1 + m) % 12 + 1
        year = start_year + (start_month - 1 + m) // 12
        month_labels.append(f"{year}\n{month}月" if month == 1 or m == 0 else f"{month}月")
    ax.set_xticklabels(month_labels, fontsize=8)

    ax.set_yticks(y_positions)
    ax.set_yticklabels(task_names, fontsize=9)
    ax.set_ylim(0.3, n_tasks + 0.7)
    ax.grid(axis="x", linestyle="--", alpha=0.4, color="gray")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_title("研究计划安排甘特图", fontsize=13, fontweight="bold", color=f"#{theme['primary']}", pad=10)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_formula_png(slide, latex: str, x, y, w, h, fontsize=16):
    formula_buf = formula_to_png(latex, fontsize=fontsize)
    slide.shapes.add_picture(formula_buf, Cm(x), Cm(y), Cm(w), Cm(h))


def _add_formula_note(slide, text: str, x, y, w, h, size=10, bold=False):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "微软雅黑"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb("333333")
    return box


def add_structured_formula_blocks(slide, formula_steps: list[dict], layout=None, theme=None):
    layout = layout or LAYOUT
    theme = theme or THEME
    left = layout["margin_left"]
    top = layout["header_h"] + 0.55
    total_w = layout["slide_w"] - layout["margin_left"] - layout["margin_right"]
    bottom_limit = layout["slide_h"] - layout["banner_h"] - 0.25
    gap = 0.28
    block_h = min(3.0, max(2.0, (bottom_limit - top - gap * max(0, len(formula_steps) - 1)) / max(1, len(formula_steps))))

    for idx, step in enumerate(formula_steps):
        y = top + idx * (block_h + gap)
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(y), Cm(total_w), Cm(block_h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb("FFFFFF")
        panel.line.color.rgb = rgb("B7B7B7")
        panel.line.width = Pt(1)

        step_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left - 0.08), Cm(y + 0.12), Cm(1.1), Cm(0.62))
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = rgb("EFEFEF")
        step_box.line.color.rgb = rgb("A6A6A6")
        step_box.line.width = Pt(0.8)
        step_tf = step_box.text_frame
        step_p = step_tf.paragraphs[0]
        step_p.alignment = PP_ALIGN.CENTER
        run = step_p.add_run()
        run.text = str(step.get("step", idx + 1))
        run.font.name = "Times New Roman"
        run.font.size = Pt(13)
        run.font.bold = True

        title = step.get("title", f"步骤 {idx + 1}")
        _add_formula_note(slide, title, left + 1.2, y + 0.08, total_w - 1.5, 0.55, size=15, bold=True)

        equation_y = y + 0.82
        equation = step.get("equation")
        if equation:
            _add_formula_png(slide, equation, left + 1.0, equation_y, total_w - 2.0, 0.88, fontsize=17)
            equation_y += 0.92
        result = step.get("result")
        if result:
            _add_formula_png(slide, result, left + 1.0, equation_y, total_w - 2.0, 0.78, fontsize=16)
            equation_y += 0.82

        note = step.get("note")
        if note:
            _add_formula_note(slide, note, left + 1.0, equation_y, total_w - 2.0, max(0.45, y + block_h - equation_y - 0.15), size=10)


def add_sectioned_formula_flow(slide, sections: list[dict], layout=None, theme=None):
    layout = layout or LAYOUT
    theme = theme or THEME
    left = layout["margin_left"]
    top = layout["header_h"] + 0.55
    section_gap = 0.38
    total_w = layout["slide_w"] - layout["margin_left"] - layout["margin_right"]
    bottom_limit = layout["slide_h"] - layout["banner_h"] - 0.25
    section_h = min(4.2, max(3.0, (bottom_limit - top - section_gap * max(0, len(sections) - 1)) / max(1, len(sections))))

    for idx, section in enumerate(sections):
        y = top + idx * (section_h + section_gap)
        title_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(left), Cm(y), Cm(total_w), Cm(0.72))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = rgb(theme["primary"])
        title_box.line.fill.background()
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.alignment = PP_ALIGN.LEFT
        run = title_p.add_run()
        run.text = section.get("title", f"模块 {idx + 1}")
        run.font.name = "微软雅黑"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = rgb("FFFFFF")

        content_y = y + 0.9
        explanation = section.get("explanation")
        if explanation:
            _add_formula_note(slide, explanation, left + 0.2, content_y, total_w * 0.42, 1.1, size=10)
        formulas = section.get("formulas", [])
        formula_area_x = left + (0.2 if not explanation else total_w * 0.45)
        formula_area_w = total_w - (formula_area_x - left) - 0.2
        if formulas:
            each_h = min(0.9, max(0.55, (section_h - 1.05) / len(formulas)))
            for f_idx, formula in enumerate(formulas):
                fy = content_y + f_idx * each_h
                _add_formula_png(slide, formula, formula_area_x, fy, formula_area_w, min(0.75, each_h - 0.05), fontsize=16)

        detail_lines = section.get("details", [])
        if detail_lines:
            detail_text = "；".join(detail_lines)
            _add_formula_note(slide, detail_text, left + 0.2, y + section_h - 0.7, total_w - 0.4, 0.45, size=9)


def make_formula_slide(
    prs,
    slide_num,
    main_title,
    paper_subtitle,
    formula_mode: str,
    formula_content: list[dict],
    bottom_banner: str,
    citations_for_this_slide=None,
    theme=None,
    layout=None,
    fonts=None,
):
    theme = theme or THEME
    layout = layout or LAYOUT
    fonts = fonts or FONTS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    if formula_mode == "modular":
        add_structured_formula_blocks(slide, formula_content, layout=layout, theme=theme)
    elif formula_mode == "sectioned":
        add_sectioned_formula_flow(slide, formula_content, layout=layout, theme=theme)
    else:
        raise ValueError(f"Unsupported formula mode: {formula_mode}")

    add_citation_footer(slide, citations_for_this_slide, theme, layout, fonts, include_banner=bool(bottom_banner))
    return slide


def make_gantt_slide(prs, slide_num, main_title, paper_subtitle, gantt_data: dict, bottom_banner: str, theme=None, layout=None, fonts=None):
    theme = theme or THEME
    layout = layout or LAYOUT
    fonts = fonts or FONTS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    gantt_buf = generate_gantt_chart(gantt_data, theme)
    chart_x = layout["margin_left"]
    chart_y = layout["header_h"] + 0.45
    chart_w = layout["slide_w"] - layout["margin_left"] - layout["margin_right"]
    chart_h = layout["slide_h"] - chart_y - 0.5
    slide.shapes.add_picture(gantt_buf, Cm(chart_x), Cm(chart_y), Cm(chart_w), Cm(chart_h))
    return slide
