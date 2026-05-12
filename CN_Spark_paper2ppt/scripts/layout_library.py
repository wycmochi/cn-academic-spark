from __future__ import annotations

import re

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

from CN_Spark_workflow.scripts.diagram_templates import (
    add_bottom_banner as workflow_add_bottom_banner,
)
from CN_Spark_workflow.scripts.diagram_templates import add_header as workflow_add_header
from CN_Spark_workflow.scripts.diagram_templates import (
    make_matrix_framework_slide as workflow_make_matrix_framework_slide,
)
from CN_Spark_workflow.scripts.diagram_templates import (
    make_mind_map_slide,
    make_network_slide,
    make_pipeline_slide as workflow_make_pipeline_slide,
    make_timeline_slide,
)


def cm(val):
    return Cm(val)


def rgb(hex_str):
    if isinstance(hex_str, (list, tuple)):
        return RGBColor(*hex_str)
    s = str(hex_str).lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


THEME = {
    "primary": "1F3864",
    "secondary_blue": "4472C4",
    "bg_gray": "F0F4FA",
    "neutral_gray": "CCCCCC",
    "accent_red": "C00000",
    "white": "FFFFFF",
}

FONTS = {
    "slide_number": ("微软雅黑", 28, True, "FFFFFF"),
    "main_title": ("微软雅黑", 20, True, "FFFFFF"),
    "paper_subtitle": ("微软雅黑", 12, False, "BDD7EE"),
    "banner_text": ("微软雅黑", 11, False, "FFFFFF"),
    "caption": ("微软雅黑", 9, False, "AAAAAA"),
}

LAYOUT = {
    "slide_w": 29.7,
    "slide_h": 21.0,
    "margin_left": 1.5,
    "margin_right": 1.5,
    "margin_top": 2.0,
    "header_h": 1.6,
    "banner_h": 1.1,
    "logo_w": 3.0,
    "logo_h": 1.0,
    "cite_h": 1.6,
}


def _resolved(theme=None, layout=None, fonts=None):
    return theme or THEME, layout or LAYOUT, fonts or FONTS


def add_header(slide, slide_number: str, main_title: str, paper_subtitle: str, theme: dict, layout: dict, fonts: dict):
    workflow_add_header(slide, slide_number, main_title, paper_subtitle, theme, layout, fonts)


def add_bottom_banner(slide, banner_text: str, theme: dict, layout: dict, fonts: dict):
    workflow_add_bottom_banner(slide, banner_text, theme, layout, fonts)


def _write_mixed_runs(paragraph, text: str, size_pt: int, color_hex: str):
    parts = re.findall(r"[\u4e00-\u9fff]+|[^\u4e00-\u9fff]+", text)
    for part in parts:
        run = paragraph.add_run()
        run.text = part
        run.font.name = "微软雅黑" if re.match(r"^[\u4e00-\u9fff]+$", part) else "Times New Roman"
        run.font.size = Pt(size_pt)
        run.font.color.rgb = rgb(color_hex)


def add_citation_footer(slide, citations_for_this_slide, theme=None, layout=None, fonts=None, include_banner=True, overflow_note=None):
    theme, layout, fonts = _resolved(theme, layout, fonts)
    citations = [c for c in (citations_for_this_slide or []) if c]
    if not citations and not overflow_note:
        return None

    cite_h = layout.get("cite_h", 1.6)
    footer_bottom = layout["slide_h"] - (layout["banner_h"] if include_banner else 0.45)
    footer_y = footer_bottom - cite_h - (0.08 if include_banner else 0)
    footer_x = layout["margin_left"]
    footer_w = layout["slide_w"] - layout["margin_left"] - layout["margin_right"]

    box = slide.shapes.add_textbox(cm(footer_x), cm(footer_y), cm(footer_w), cm(cite_h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    visible = citations[:3]
    for idx, cite in enumerate(visible):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        _write_mixed_runs(p, cite, 8, "888888")

    if overflow_note:
        p = tf.add_paragraph() if visible else tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        _write_mixed_runs(p, overflow_note, 8, "888888")
    elif len(citations) > 3:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        _write_mixed_runs(p, "其余文献详见参考文献页", 8, "888888")
    return box


def make_cover_slide(prs, title: str, subtitle: str, author: str, affiliation: str, date: str, theme=None, layout=None, fonts=None):
    theme, layout, fonts = _resolved(theme, layout, fonts)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(0), cm(0), cm(layout["slide_w"]), cm(layout["slide_h"]))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(theme["primary"])
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(cm(2.7), cm(4.2), cm(layout["slide_w"] - 5.4), cm(2.5))
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "微软雅黑"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = rgb(theme["white"])

    meta_lines = [line for line in [subtitle, author, affiliation, date] if line]
    meta_box = slide.shapes.add_textbox(cm(4.2), cm(8.0), cm(layout["slide_w"] - 8.4), cm(5.0))
    meta_tf = meta_box.text_frame
    meta_tf.word_wrap = True
    for idx, line in enumerate(meta_lines):
        p = meta_tf.paragraphs[0] if idx == 0 else meta_tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = "微软雅黑"
        run.font.size = Pt(14 if idx == 0 else 12)
        run.font.color.rgb = rgb("D9E2F3")
    return slide


def make_pipeline_slide(prs, slide_num: str, main_title: str, paper_subtitle: str, stages: list[dict], arrows: list[dict], bottom_banner: str, theme=None, layout=None, fonts=None):
    theme, layout, fonts = _resolved(theme, layout, fonts)
    return workflow_make_pipeline_slide(prs, slide_num, main_title, paper_subtitle, stages, arrows, bottom_banner, theme, layout, fonts)


def make_matrix_framework_slide(prs, slide_num: str, main_title: str, paper_subtitle: str, row_groups: list[dict], mid_nodes: list[dict], right_texts: list[str], left_to_mid_edges: list[tuple], bottom_banner: str, theme=None, layout=None, fonts=None):
    theme, layout, fonts = _resolved(theme, layout, fonts)
    return workflow_make_matrix_framework_slide(
        prs,
        slide_num,
        main_title,
        paper_subtitle,
        row_groups,
        mid_nodes,
        right_texts,
        left_to_mid_edges,
        bottom_banner,
        theme,
        layout,
        fonts,
    )


def make_conceptual_framework_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    framework_type: str,
    framework_data: dict,
    bottom_banner: str,
    citations_for_this_slide=None,
    theme=None,
    layout=None,
    fonts=None,
):
    theme, layout, fonts = _resolved(theme, layout, fonts)
    framework_type = framework_type.lower()
    if framework_type in {"mind_map", "mindmap", "radial"}:
        slide = make_mind_map_slide(
            prs,
            slide_num,
            main_title,
            paper_subtitle,
            framework_data.get("center_node", {}),
            framework_data.get("level1_nodes", []),
            framework_data.get("level2_groups", {}),
            framework_data.get("leaf_citations", {}),
            bottom_banner,
            theme,
            layout,
            fonts,
        )
    elif framework_type in {"network", "concept_network"}:
        slide = make_network_slide(
            prs,
            slide_num,
            main_title,
            paper_subtitle,
            framework_data.get("nodes", []),
            framework_data.get("edges", []),
            bottom_banner,
            theme,
            layout,
            fonts,
        )
    elif framework_type in {"timeline", "evolution_timeline"}:
        slide = make_timeline_slide(
            prs,
            slide_num,
            main_title,
            paper_subtitle,
            framework_data.get("timeline_points", []),
            framework_data.get("relation_edges", []),
            bottom_banner,
            theme,
            layout,
            fonts,
        )
    else:
        raise ValueError(f"Unsupported conceptual framework type: {framework_type}")

    add_citation_footer(slide, citations_for_this_slide, theme, layout, fonts, include_banner=bool(bottom_banner))
    return slide
