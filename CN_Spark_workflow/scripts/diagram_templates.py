from __future__ import annotations

from math import cos, pi, sin

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


def cm(val):
    return Cm(val)


def rgb(hex_str):
    if isinstance(hex_str, (list, tuple)):
        return RGBColor(*hex_str)
    s = str(hex_str).lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts):
    header_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        cm(0),
        cm(0),
        cm(layout["slide_w"]),
        cm(layout["header_h"]),
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = rgb(theme["primary"])
    header_bg.line.fill.background()

    num_box = slide.shapes.add_textbox(cm(0.3), cm(0.08), cm(1.9), cm(1.1))
    num_tf = num_box.text_frame
    num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    num_p = num_tf.paragraphs[0]
    num_p.alignment = PP_ALIGN.LEFT
    run = num_p.add_run()
    run.text = str(slide_num or "")
    run.font.name = fonts["slide_number"][0]
    run.font.size = Pt(fonts["slide_number"][1])
    run.font.bold = fonts["slide_number"][2]
    run.font.color.rgb = rgb(fonts["slide_number"][3])

    title_box = slide.shapes.add_textbox(
        cm(2.4),
        cm(0.15),
        cm(layout["slide_w"] - 6.3),
        cm(0.7),
    )
    title_tf = title_box.text_frame
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    run = title_p.add_run()
    run.text = main_title or ""
    run.font.name = fonts["main_title"][0]
    run.font.size = Pt(fonts["main_title"][1])
    run.font.bold = fonts["main_title"][2]
    run.font.color.rgb = rgb(fonts["main_title"][3])

    if paper_subtitle:
        subtitle_box = slide.shapes.add_textbox(
            cm(2.4),
            cm(0.84),
            cm(layout["slide_w"] - 6.3),
            cm(0.45),
        )
        subtitle_tf = subtitle_box.text_frame
        subtitle_p = subtitle_tf.paragraphs[0]
        subtitle_p.alignment = PP_ALIGN.LEFT
        run = subtitle_p.add_run()
        run.text = paper_subtitle
        run.font.name = fonts["paper_subtitle"][0]
        run.font.size = Pt(fonts["paper_subtitle"][1])
        run.font.bold = fonts["paper_subtitle"][2]
        run.font.color.rgb = rgb(fonts["paper_subtitle"][3])


def add_bottom_banner(slide, banner_text, theme, layout, fonts):
    if not banner_text:
        return
    banner_y = layout["slide_h"] - layout["banner_h"]
    banner_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        cm(0),
        cm(banner_y),
        cm(layout["slide_w"]),
        cm(layout["banner_h"]),
    )
    banner_bg.fill.solid()
    banner_bg.fill.fore_color.rgb = rgb(theme["primary"])
    banner_bg.line.fill.background()

    text_box = slide.shapes.add_textbox(
        cm(1.4),
        cm(banner_y + 0.08),
        cm(layout["slide_w"] - 2.8),
        cm(layout["banner_h"] - 0.16),
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = banner_text
    run.font.name = fonts["banner_text"][0]
    run.font.size = Pt(fonts["banner_text"][1])
    run.font.bold = fonts["banner_text"][2]
    run.font.color.rgb = rgb(fonts["banner_text"][3])


def _content_bounds(layout, include_banner=True):
    left = layout["margin_left"]
    top = layout["header_h"] + 0.45
    right = layout["slide_w"] - layout["margin_right"]
    bottom = layout["slide_h"] - (layout["banner_h"] + 0.35 if include_banner else 0.45)
    return left, top, right, bottom


def _add_text(slide, text, x, y, w, h, font_name, size_pt, bold=False, color_hex="000000", align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_hex)
    return box


def _add_node(slide, x, y, w, h, title, theme, fill_hex, text_hex, subtitle=None, citation=None, shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    node = slide.shapes.add_shape(shape, cm(x), cm(y), cm(w), cm(h))
    node.fill.solid()
    node.fill.fore_color.rgb = rgb(fill_hex)
    node.line.color.rgb = rgb(theme["primary"])
    node.line.width = Pt(1)

    tf = node.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    title_run = p.add_run()
    title_run.text = title
    title_run.font.name = "微软雅黑"
    title_run.font.size = Pt(13)
    title_run.font.bold = True
    title_run.font.color.rgb = rgb(text_hex)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        subtitle_run = p2.add_run()
        subtitle_run.text = subtitle
        subtitle_run.font.name = "微软雅黑"
        subtitle_run.font.size = Pt(9)
        subtitle_run.font.color.rgb = rgb(text_hex)

    if citation:
        _add_text(slide, citation, x + w - 0.95, y + h - 0.38, 0.85, 0.25, "Times New Roman", 8, False, "888888", PP_ALIGN.RIGHT)
    return node


def _connect(slide, x1, y1, x2, y2, color_hex="555555", width_pt=1.8, elbow=False):
    conn_type = MSO_CONNECTOR_TYPE.ELBOW if elbow else MSO_CONNECTOR_TYPE.STRAIGHT
    line = slide.shapes.add_connector(conn_type, cm(x1), cm(y1), cm(x2), cm(y2))
    line.line.color.rgb = rgb(color_hex)
    line.line.width = Pt(width_pt)
    return line


def make_matrix_framework_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    row_groups: list[dict],
    mid_nodes: list[dict],
    right_texts: list[str],
    left_to_mid_edges: list[tuple],
    bottom_banner: str,
    theme: dict,
    layout: dict,
    fonts: dict,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    left, top, right, bottom = _content_bounds(layout)
    col_gap = 0.5
    total_w = right - left
    left_w = total_w * 0.27
    mid_w = total_w * 0.34
    right_w = total_w - left_w - mid_w - col_gap * 2
    header_h = 0.8
    body_h = bottom - top - header_h - 0.25

    columns = [
        (left, left_w, "研究维度"),
        (left + left_w + col_gap, mid_w, "核心模块"),
        (left + left_w + mid_w + col_gap * 2, right_w, "输出 / 结论"),
    ]
    for x, w, title in columns:
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(x), cm(top), cm(w), cm(header_h))
        head.fill.solid()
        head.fill.fore_color.rgb = rgb(theme["primary"])
        head.line.fill.background()
        _add_text(slide, title, x, top + 0.08, w, header_h - 0.16, "微软雅黑", 12, True, "FFFFFF")

    def _stack_cards(items, x, w, start_y):
        if not items:
            return []
        gap = 0.22
        card_h = min(1.45, (body_h - gap * (len(items) - 1)) / len(items))
        nodes = []
        for idx, item in enumerate(items):
            y = start_y + idx * (card_h + gap)
            title = item["title"] if isinstance(item, dict) else str(item)
            subtitle = item.get("subtitle") if isinstance(item, dict) else None
            citation = item.get("citation") if isinstance(item, dict) else None
            nodes.append(
                _add_node(
                    slide,
                    x,
                    y,
                    w,
                    card_h,
                    title,
                    theme,
                    theme["bg_gray"],
                    theme["primary"],
                    subtitle=subtitle,
                    citation=citation,
                )
            )
        return nodes

    start_y = top + header_h + 0.22
    left_nodes = _stack_cards(row_groups, left, left_w, start_y)
    mid_shapes = _stack_cards(mid_nodes, left + left_w + col_gap, mid_w, start_y)
    right_items = [{"title": t} if not isinstance(t, dict) else t for t in right_texts]
    right_nodes = _stack_cards(right_items, left + left_w + mid_w + col_gap * 2, right_w, start_y)

    all_mid = mid_shapes or []
    for left_idx, mid_idx in left_to_mid_edges or []:
        if 0 <= left_idx < len(left_nodes) and 0 <= mid_idx < len(all_mid):
            l = left_nodes[left_idx]
            m = all_mid[mid_idx]
            _connect(
                slide,
                l.left.cm + l.width.cm,
                l.top.cm + l.height.cm / 2,
                m.left.cm,
                m.top.cm + m.height.cm / 2,
            )

    pair_count = min(len(all_mid), len(right_nodes))
    for idx in range(pair_count):
        m = all_mid[idx]
        r = right_nodes[idx]
        _connect(
            slide,
            m.left.cm + m.width.cm,
            m.top.cm + m.height.cm / 2,
            r.left.cm,
            r.top.cm + r.height.cm / 2,
        )
    return slide


def make_pipeline_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    stages: list[dict],
    arrows: list[dict],
    bottom_banner: str,
    theme: dict,
    layout: dict,
    fonts: dict,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    left, top, right, bottom = _content_bounds(layout)
    content_w = right - left
    node_h = min(3.0, bottom - top - 0.6)
    node_y = top + ((bottom - top) - node_h) / 2
    stage_count = max(1, len(stages))
    gap = min(0.7, content_w * 0.03)
    node_w = (content_w - gap * (stage_count - 1)) / stage_count
    nodes = []

    for idx, stage in enumerate(stages):
        x = left + idx * (node_w + gap)
        shape = MSO_AUTO_SHAPE_TYPE.CHEVRON if stage.get("kind") == "decision" else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        nodes.append(
            _add_node(
                slide,
                x,
                node_y,
                node_w,
                node_h,
                stage.get("title", ""),
                theme,
                stage.get("fill", theme["secondary_blue"]),
                stage.get("text_color", theme["white"]),
                subtitle=stage.get("subtitle"),
                citation=stage.get("citation"),
                shape=shape,
            )
        )

    for idx in range(len(nodes) - 1):
        a = nodes[idx]
        b = nodes[idx + 1]
        _connect(
            slide,
            a.left.cm + a.width.cm,
            a.top.cm + a.height.cm / 2,
            b.left.cm,
            b.top.cm + b.height.cm / 2,
        )

    for arrow in arrows or []:
        src = arrow.get("from")
        dst = arrow.get("to")
        if src is None or dst is None or src >= len(nodes) or dst >= len(nodes):
            continue
        a = nodes[src]
        b = nodes[dst]
        _connect(
            slide,
            a.left.cm + a.width.cm / 2,
            a.top.cm + a.height.cm,
            b.left.cm + b.width.cm / 2,
            b.top.cm,
            color_hex=arrow.get("color", theme["accent_red"]),
            width_pt=1.2,
            elbow=True,
        )
        if arrow.get("label"):
            label_x = (a.left.cm + b.left.cm) / 2
            label_y = node_y + node_h + 0.2
            _add_text(slide, arrow["label"], label_x, label_y, 2.3, 0.45, "微软雅黑", 9, False, "666666")
    return slide


def make_mind_map_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    center_node: dict,
    level1_nodes: list[dict],
    level2_groups: dict,
    leaf_citations: dict,
    bottom_banner: str,
    theme: dict,
    layout: dict,
    fonts: dict,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    left, top, right, bottom = _content_bounds(layout)
    cx = (left + right) / 2
    cy = (top + bottom) / 2
    center_w, center_h = 4.2, 1.5
    _add_node(
        slide,
        cx - center_w / 2,
        cy - center_h / 2,
        center_w,
        center_h,
        center_node.get("title", ""),
        theme,
        theme["primary"],
        theme["white"],
        subtitle=center_node.get("subtitle"),
        citation=center_node.get("citation"),
    )

    branch_count = max(1, len(level1_nodes))
    radius_1 = min((right - left) * 0.25, (bottom - top) * 0.28)
    radius_2 = radius_1 + 2.6
    level1_shapes = {}
    center_id = center_node.get("id", center_node.get("title", "center"))

    for idx, node in enumerate(level1_nodes):
        angle = (-pi / 2) + 2 * pi * idx / branch_count
        w, h = 3.2, 1.15
        x = cx + cos(angle) * radius_1 - w / 2
        y = cy + sin(angle) * radius_1 - h / 2
        node_id = node.get("id", node.get("title", f"l1_{idx}"))
        level1_shapes[node_id] = (x, y, w, h, angle)
        _add_node(slide, x, y, w, h, node.get("title", ""), theme, theme["secondary_blue"], theme["white"], subtitle=node.get("subtitle"))
        _connect(slide, cx, cy, x + w / 2, y + h / 2, elbow=False)

    for parent_id, children in (level2_groups or {}).items():
        parent = level1_shapes.get(parent_id)
        if not parent:
            continue
        px, py, pw, ph, angle = parent
        child_count = max(1, len(children))
        spread = pi / 6 if child_count > 1 else 0
        for idx, child in enumerate(children):
            offset = -spread / 2 + (spread * idx / max(1, child_count - 1)) if child_count > 1 else 0
            child_angle = angle + offset
            w, h = 2.55, 0.95
            x = cx + cos(child_angle) * radius_2 - w / 2
            y = cy + sin(child_angle) * radius_2 - h / 2
            child_id = child.get("id", child.get("title", f"{parent_id}_{idx}"))
            citation = leaf_citations.get(child_id) or child.get("citation")
            _add_node(slide, x, y, w, h, child.get("title", ""), theme, "DCE6F2", theme["primary"], citation=citation)
            _connect(slide, px + pw / 2, py + ph / 2, x + w / 2, y + h / 2, elbow=False)

    if center_node.get("citation") or leaf_citations.get(center_id):
        _add_text(slide, center_node.get("citation") or leaf_citations.get(center_id), cx + 1.4, cy + 0.65, 0.9, 0.25, "Times New Roman", 8, False, "888888")
    return slide


def make_network_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    nodes: list[dict],
    edges: list[dict],
    bottom_banner: str,
    theme: dict,
    layout: dict,
    fonts: dict,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    left, top, right, bottom = _content_bounds(layout)
    width = right - left
    height = bottom - top
    radius = min(width, height) * 0.34
    cx = left + width * 0.43
    cy = top + height * 0.5
    node_count = max(1, len(nodes))
    placed = {}

    for idx, node in enumerate(nodes):
        angle = (-pi / 2) + 2 * pi * idx / node_count
        importance = max(1.0, float(node.get("size", 1.0)))
        w = min(3.4, 2.1 + importance * 0.2)
        h = 0.95 + min(0.45, importance * 0.05)
        x = cx + cos(angle) * radius - w / 2
        y = cy + sin(angle) * radius - h / 2
        node_id = node.get("id", idx)
        placed[node_id] = (x, y, w, h)
        _add_node(
            slide,
            x,
            y,
            w,
            h,
            node.get("title", ""),
            theme,
            node.get("fill", "DCE6F2"),
            node.get("text_color", theme["primary"]),
            citation=node.get("citation"),
        )

    for edge in edges or []:
        src = placed.get(edge.get("from"))
        dst = placed.get(edge.get("to"))
        if not src or not dst:
            continue
        sx, sy, sw, sh = src
        dx, dy, dw, dh = dst
        color = theme["accent_red"] if edge.get("kind") == "conflict" else theme["secondary_blue"]
        width_pt = 1.2 if edge.get("kind") == "support" else 1.8
        _connect(
            slide,
            sx + sw / 2,
            sy + sh / 2,
            dx + dw / 2,
            dy + dh / 2,
            color_hex=color,
            width_pt=width_pt,
        )
        if edge.get("label"):
            mid_x = (sx + dx) / 2
            mid_y = (sy + dy) / 2
            _add_text(slide, edge["label"], mid_x, mid_y, 1.8, 0.35, "微软雅黑", 8, False, "666666")

    legend_x = left + width * 0.78
    legend_y = top + 0.4
    _add_text(slide, "关系图例", legend_x, legend_y, 2.5, 0.5, "微软雅黑", 11, True, theme["primary"], PP_ALIGN.LEFT)
    legend_items = [
        ("support", theme["secondary_blue"], "依赖 / 支撑"),
        ("conflict", theme["accent_red"], "争议 / 冲突"),
    ]
    for idx, (_, color_hex, label) in enumerate(legend_items):
        y = legend_y + 0.7 + idx * 0.55
        _connect(slide, legend_x, y + 0.16, legend_x + 0.8, y + 0.16, color_hex=color_hex, width_pt=1.8)
        _add_text(slide, label, legend_x + 0.95, y, 2.1, 0.35, "微软雅黑", 9, False, "666666", PP_ALIGN.LEFT)
    return slide


def make_timeline_slide(
    prs,
    slide_num: str,
    main_title: str,
    paper_subtitle: str,
    timeline_points: list[dict],
    relation_edges: list[dict],
    bottom_banner: str,
    theme: dict,
    layout: dict,
    fonts: dict,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, slide_num, main_title, paper_subtitle, theme, layout, fonts)
    add_bottom_banner(slide, bottom_banner, theme, layout, fonts)

    left, top, right, bottom = _content_bounds(layout)
    axis_y = top + 1.0
    _connect(slide, left, axis_y, right, axis_y, color_hex=theme["primary"], width_pt=2.4)
    if not timeline_points:
        return slide

    years = [int(p.get("year", 0)) for p in timeline_points]
    min_year, max_year = min(years), max(years)
    span = max(1, max_year - min_year)
    placed = {}

    for idx, point in enumerate(sorted(timeline_points, key=lambda p: int(p.get("year", 0)))):
        year = int(point.get("year", min_year))
        ratio = (year - min_year) / span
        x_center = left + ratio * (right - left)
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cm(x_center - 0.03), cm(axis_y - 0.18), cm(0.06), cm(0.36))
        tick.fill.solid()
        tick.fill.fore_color.rgb = rgb(theme["primary"])
        tick.line.fill.background()
        _add_text(slide, str(year), x_center - 0.65, axis_y - 0.85, 1.3, 0.38, "Times New Roman", 10, True, theme["primary"])

        band = idx % 2
        w, h = 2.7, 1.05
        y = axis_y + 0.55 + band * 2.2
        fill_hex = "D9E2F3" if ratio < 0.33 else theme["secondary_blue"] if ratio < 0.66 else "E06666"
        text_hex = theme["primary"] if fill_hex == "D9E2F3" else theme["white"]
        x = min(max(left, x_center - w / 2), right - w)
        placed[point.get("id", idx)] = (x, y, w, h)
        _add_node(slide, x, y, w, h, point.get("title", ""), theme, fill_hex, text_hex, subtitle=point.get("subtitle"), citation=point.get("citation"))
        _connect(slide, x_center, axis_y, x + w / 2, y, color_hex="888888", width_pt=1.2)

    for edge in relation_edges or []:
        src = placed.get(edge.get("from"))
        dst = placed.get(edge.get("to"))
        if not src or not dst:
            continue
        sx, sy, sw, sh = src
        dx, dy, dw, dh = dst
        kind = edge.get("kind", "derive")
        color = theme["accent_red"] if kind == "challenge" else "555555"
        _connect(
            slide,
            sx + sw / 2,
            sy + sh / 2,
            dx + dw / 2,
            dy + dh / 2,
            color_hex=color,
            width_pt=1.3,
        )
    return slide
